#!/usr/bin/env python3
"""Generate PowerPoint: Giai đoạn Tư vấn (Consult) — Service Lifecycle CRM PTT."""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from crm_svc_tasks import SERVICE_LABELS
from crm_svc_workflow_steps import SERVICE_WORKFLOW_STEPS

OUTPUT = ROOT / "docs" / "Consult_Stage_Service_Delivery.pptx"
LOGO = ROOT / "static" / "images" / "ptt-logo.png"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x63, 0x66, 0xF1)
TEAL = RGBColor(0x0D, 0x94, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT_LIGHT = RGBColor(0xC7, 0xD2, 0xFE)
WARN = RGBColor(0xB4, 0x53, 0x09)


def _logo_exists() -> bool:
    return LOGO.is_file()


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _header(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    if _logo_exists():
        slide.shapes.add_picture(str(LOGO), Inches(8.35), Inches(0.08), height=Inches(0.88))
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(7.7), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(7.7), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = ACCENT_LIGHT


def _bullets(slide, items: list[str], top=1.15, size=13, left=0.45, width=9.1, height=6.0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        lvl = 1 if item.startswith("  ") else 0
        p.level = lvl
        p.font.size = Pt(size - 1 if lvl else size)
        p.font.color.rgb = GRAY if lvl else DARK_TEXT
        p.space_after = Pt(4)


def _table_slide(
    prs,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    subtitle: str = "",
    col_widths: list[float] | None = None,
):
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, title, subtitle)
    nr, nc = len(rows) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(nr, nc, Inches(0.25), Inches(1.12), Inches(9.5), Inches(5.85))
    tbl = tbl_shape.table
    if col_widths and len(col_widths) == nc:
        total = sum(col_widths)
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(9.5 * w / total)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK_TEXT
            if i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT_BG
    return slide


def _title_slide(prs: Presentation):
    slide = prs.slides.add_slide(_blank(prs))
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    if _logo_exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.1), Inches(0.85), height=Inches(1.35))
    accent = slide.shapes.add_shape(1, 0, Inches(2.35), Inches(10), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    for y, text, sz, bold, color in [
        (2.65, "Giai đoạn Tư vấn (Consult)", 34, True, WHITE),
        (3.35, "Service Lifecycle CRM — PTT Advertising", 18, False, ACCENT_LIGHT),
        (4.05, "Brief Lead · Prefill · Gate funnel · AI · Proposal bridge", 14, False, GRAY),
        (4.65, "Spec v1.0 · 2026-06-30 · docs/specs/2026-06-30-consult-stage-system-design.md", 11, False, GRAY),
    ]:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER


def _closing_slide(prs: Presentation):
    slide = prs.slides.add_slide(_blank(prs))
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    if _logo_exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.25), Inches(2.0), height=Inches(1.2))
    for y, text, sz in [
        (3.45, "Cảm ơn", 40),
        (4.35, "PTT Advertising Solutions", 16),
        (4.85, "CRM Service Delivery · Consult Stage v1", 12),
    ]:
        tb = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = sz >= 20
        p.font.color.rgb = WHITE if sz >= 20 else GRAY
        p.alignment = PP_ALIGN.CENTER


def _consult_rows_from_code() -> list[list[str]]:
    rows: list[list[str]] = []
    for i, slug in enumerate(sorted(SERVICE_WORKFLOW_STEPS.keys()), 1):
        steps = SERVICE_WORKFLOW_STEPS[slug].get("consult") or []
        if not steps:
            continue
        step = steps[0]
        fields = ", ".join(f["key"] for f in (step.get("form_fields") or []))
        rows.append([
            str(i),
            SERVICE_LABELS.get(slug, slug),
            step.get("title", "")[:45],
            fields[:80],
        ])
    return rows


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Mục lục")
    _bullets(slide, [
        "I. Vấn đề & hiện trạng hệ thống",
        "II. Nghiệp vụ Consult trong vòng đời 7 giai đoạn",
        "III. Input từ Lead · BANT+ · Go/Nurture/No-Go",
        "IV. Giải pháp: Consult Brief · Prefill · Gate · AI · Proposal",
        "V. Consult task theo 12 dịch vụ (từ CRM template)",
        "VI. RACI · SLA · KPI funnel",
        "VII. Kế hoạch triển khai C0–C6 · Go-live",
    ], top=1.15, size=14)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Vấn đề hiện tại", "Khảo sát codebase PTTADS · 2026-06-30")
    _bullets(slide, [
        "Consult = stage generic trong workflow — không có module bridge riêng",
        "Lead Intake mạnh ở stage Lead nhưng DỪNG ở merge task Lead",
        "Form Consult mở ra TRỐNG — AM nhập lại niche/budget/pain",
        "AI consult_analysis thiếu BANT/decision/Intake context",
        "Proposal AI không đọc audit Consult",
        "no_go vẫn chuyển Consult nếu tick xong task Lead",
        "",
        "Hệ quả: lãng phí effort, báo giá thiếu scope, không đo funnel Go→Consult",
    ], size=13)

    _table_slide(
        prs,
        "7 giai đoạn Service Lifecycle",
        ["#", "Stage", "Mục tiêu", "Sync Lead care"],
        [
            ["1", "Lead", "Qualify sơ bộ · Intake phone/B", "first_contact"],
            ["2", "Tư vấn", "Audit/discovery sâu", "qualify"],
            ["3", "Báo giá", "Proposal + KPI cam kết", "advise"],
            ["4", "Onboarding", "Kickoff · access", "closing"],
            ["5", "Triển khai", "Deliver hàng tháng/sprint", "post_sale"],
            ["6", "Nghiệm thu", "Handover KPI", "post_sale"],
            ["7", "Chăm sóc", "Retain · upsell", "post_sale"],
        ],
        subtitle="Chỉ được chuyển tuần tự — hoàn thành 100% task stage hiện tại",
        col_widths=[0.4, 1.0, 3.5, 1.2],
    )

    _table_slide(
        prs,
        "BANT+ & Quyết định (Lead Intake)",
        ["Điểm", "Decision", "Consult", "Hành động"],
        [
            ["≥ 24", "Go", "Consult đầy đủ", "Audit → Proposal ≤48h"],
            ["18–23", "Nurture", "Consult nhẹ / follow-up", "Drip 7–14 ngày"],
            ["< 18", "No-Go", "Không Consult sâu", "Đóng lifecycle / chuyển DV"],
            ["≥2 red flags", "No-Go (gợi ý)", "Cảnh báo", "Director override nếu Go"],
        ],
        subtitle="6 tiêu chí × 1–5: budget, authority, need, timeline, fit, history",
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Input bắt buộc từ Lead → Consult")
    _bullets(slide, [
        "Từ task Lead (form_data): niche, budget, need, domain… theo từng dịch vụ",
        "Từ Lead Intake session: BANT, decision, answers, stakeholders, 3 cam kết",
        "Từ AI Intake: ai_summary, red flags, recommended_next_step",
        "Lịch: next_meeting_at, proposal_date",
        "",
        "Gate trước khi chuyển Lead → Consult:",
        "  ✓ 100% task Lead done (engine hiện có)",
        "  ✓ decision set · cảnh báo no_go / nurture (triển khai C3)",
        "  ✓ Khuyến nghị: có Intake in_person nếu decision=go",
    ], size=12)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Luồng mục tiêu (Target state)", "Sau triển khai Consult Bridge v1")
    _bullets(slide, [
        "1. Intake phone (PHẦN A) → decision + BANT",
        "2. Go → hẹn Intake gặp (PHẦN B) → auto ✓ task Lead nếu BANT≥24",
        "3. AM: Chuyển → Tư vấn → prefill Consult + hiện Consult Brief",
        "4. Audit theo task Consult (12 DV) + AI consult_analysis (full context)",
        "5. Tick ✓ Consult → Chuyển → Báo giá (Proposal đọc Consult output)",
        "",
        "Khác biệt Lead vs Consult:",
        "  Lead = qualify + BANT + cam kết (Intake form)",
        "  Consult = audit chuyên môn + scope cho báo giá (task CRM)",
        "  Consult miễn phí — PTT không thu phí tư vấn; doanh thu chỉ sau ký HĐ",
    ], size=12)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Giải pháp C1 — Consult Brief Panel", "UI trên /crm/service-delivery/<lifecycle_id>")
    _bullets(slide, [
        "Hiển thị khi lifecycle.stage = consult",
        "Tổng hợp: decision, BANT, temperature, pain, budget, domain",
        "Danh sách Intake sessions (gọi/gặp) + AI summary",
        "Stakeholders · 3 cam kết KH · proposal_date",
        "Cảnh báo: chưa có PHẦN B in_person",
        "Nút: Prefill form · Mở Intake · Xem task Lead",
        "",
        "API: GET /api/crm/service-lifecycle/<id>/consult-brief",
        "Module: crm_svc_consult_bridge.get_consult_brief()",
    ], size=12)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Giải pháp C2 — Prefill Consult task")
    _bullets(slide, [
        "Hook: advance_stage(..., 'consult') + POST consult-prefill",
        "Chỉ điền field TRỐNG — không ghi đè AM đã sửa",
        "",
        "Mapping mẫu (SEO Tổng thể):",
        "  Lead need → Consult current_status (Pain: …)",
        "  Intake answers → target_keywords, top_competitors",
        "  Lead niche/budget → append notes",
        "",
        "get_crm_field_map(slug) trong crm_lead_intake_definitions.py",
    ], size=12)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Giải pháp C3 — Gate & automation")
    _bullets(slide, [
        "validate_consult_advance(): ok | warn | block",
        "  block: decision=no_go (trừ Director override)",
        "  warn: nurture hoặc BANT 18–23",
        "on_intake_completed():",
        "  in_person + go + BANT≥24 → auto ✓ task Lead",
        "  phone + go → reminder Hẹn PHẦN B",
        "UI: banner Sẵn sàng chuyển Tư vấn / ẩn Lead Intake khi stage≠lead",
    ], size=12)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Giải pháp C4–C5 — AI & Proposal")
    _bullets(slide, [
        "C4: build_ai_context_for_consult() merge brief vào AI assist",
        "  Prompt consult_analysis + BANT + intake excerpt + lead form_data",
        "",
        "C5: crm_proposal.get_customer_context đọc task Consult",
        "  form_data + ai_output → Proposal AI",
        "Nút workflow: Tạo Proposal từ Consult",
    ], size=12)

    _table_slide(
        prs,
        "Consult task — 12 dịch vụ (CRM template)",
        ["#", "Dịch vụ", "Task Consult", "Form fields (keys)"],
        _consult_rows_from_code(),
        subtitle="Seed từ crm_svc_workflow_steps.py · 1 task / dịch vụ",
        col_widths=[0.35, 1.4, 2.0, 2.5],
    )

    _table_slide(
        prs,
        "RACI & SLA — Giai đoạn Consult",
        ["Hoạt động", "AM", "SP", "SLA"],
        [
            ["Đọc Consult Brief", "R", "I", "Ngay khi vào Consult"],
            ["Audit / discovery", "R", "R", "3–7 ngày"],
            ["Điền form + AI assist", "R", "C", "Trong buổi tư vấn"],
            ["Chuyển Proposal", "R", "I", "≤48h sau meeting"],
            ["Override No-Go", "R", "—", "Director duyệt"],
        ],
        col_widths=[2.5, 0.6, 0.6, 2.0],
    )

    _table_slide(
        prs,
        "KPI funnel (triển khai C6)",
        ["KPI", "Công thức", "Ngưỡng"],
        [
            ["Go → Consult", "stage≥consult & intake go / total go", "≥35%"],
            ["Consult → Proposal ≤7d", "proposal within 7d of consult", "≥60%"],
            ["Consult task done", "consult tasks ✓ / entered consult", "≥80%"],
            ["BANT avg won", "avg bant lifecycle won", "≥22"],
            ["in_person trước Consult", "in_person before consult / go", "≥80%"],
        ],
    )

    _table_slide(
        prs,
        "Kế hoạch triển khai",
        ["Phase", "Thời gian", "Deliverable"],
        [
            ["C0", "2 ngày", "SOP + PPT đào tạo (file này)"],
            ["C1", "4 ngày", "Consult Brief panel + API"],
            ["C2", "3 ngày", "Prefill Consult + field map"],
            ["C3", "4 ngày", "Gate no_go/nurture + auto ✓ Lead"],
            ["C4", "3 ngày", "AI context enrich"],
            ["C5", "4 ngày", "Proposal bridge"],
            ["C6", "2 ngày", "KPI funnel dashboard"],
        ],
        subtitle="Tổng ~3–4 tuần · Spec: docs/specs/2026-06-30-consult-stage-system-design.md",
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Checklist AM — Khi vào giai đoạn Tư vấn")
    _bullets(slide, [
        "□ Mở Consult Brief — xác nhận decision & BANT",
        "□ Nếu Go mà chưa gặp PHẦN B → hẹn Intake in_person",
        "□ Thu tài liệu L2 (GSC, GBP, Ads account…) theo dịch vụ",
        "□ Kiểm tra prefill form Consult — bổ sung field trống",
        "□ Thực hiện audit/discovery theo task title",
        "□ Chạy AI Hỗ trợ → review → chỉnh notes",
        "□ Tick ✓ task Consult",
        "□ Chuyển → Báo giá trong 48h (Proposal prefill từ Consult)",
    ], size=13)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "CRM hôm nay vs sau C1–C6", "Phase C0 = tài liệu · Code bắt đầu từ C1")
    _bullets(slide, [
        "ĐÃ CÓ:",
        "  Workflow 7 stage · task Consult template 12 DV · Lead Intake web",
        "  Chuyển Lead→Consult (tick ✓ task Lead) · AI assist cơ bản",
        "",
        "CHƯA CÓ (roadmap):",
        "  Consult Brief panel · Prefill form Consult · Gate no_go/nurture",
        "  AI đọc full Intake · Proposal prefill từ Consult",
        "",
        "AM TẠM THỜI: mở tab Lead + Intake song song khi làm Consult",
        "SOP: docs/runbooks/consult-stage-am-sop.md",
    ], size=12)

    _table_slide(
        prs,
        "Lead vs Consult — Không nhầm lẫn",
        ["", "Lead", "Consult"],
        [
            ["Công cụ", "Lead Intake PHẦN A/B + task Lead", "Task Consult workflow"],
            ["Thời lượng", "15–25p gọi · 45–60p gặp", "Audit 3–7 ngày"],
            ["Mục tiêu", "BANT+ · Go/Nurture/No-Go", "Audit · scope Proposal"],
            ["Output", "decision · cam kết KH", "form_data · AI analysis"],
            ["Doanh thu PTT", "Không", "Không — miễn phí pre-sales"],
            ["Stage care", "first_contact", "qualify"],
        ],
        col_widths=[1.2, 2.0, 2.0],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Tài liệu Phase C0 (repo)")
    _bullets(slide, [
        "docs/Consult_Stage_Service_Delivery.pptx — slide deck này",
        "docs/runbooks/consult-stage-am-sop.md — SOP AM",
        "docs/runbooks/consult-stage-training-guide.md — kịch bản 45p",
        "docs/runbooks/consult-stage-service-tasks.md — 12 task Consult",
        "docs/runbooks/consult-stage-bant-signoff.md — Director ký ngưỡng",
        "docs/crm/README.md — mục lục CRM",
        "",
        "Tạo lại: python3 scripts/generate_consult_stage_pptx.py",
        "         python3 scripts/generate_consult_runbook_appendix.py",
    ], size=12)

    slide = prs.slides.add_slide(_blank(prs))
    _bullets(slide, [
        "□ Training 45 phút với slide deck này",
        "□ Pilot 2 lifecycle: SEO Tổng thể + Facebook Ads",
        "□ Director sign-off gate No-Go / Nurture",
        "□ Monitor 2 tuần: prefill rate, consult→proposal time",
        "□ Tiêu chí nghiệm thu spec §14 — 8 checkbox",
    ], size=13)

    _closing_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out}")
