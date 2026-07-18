#!/usr/bin/env python3
"""Generate PowerPoint: Checklist tiếp nhận Lead — 12 dịch vụ PTT (chi tiết + logo)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "Checklist_Tiep_Nhan_Lead_12_Dich_Vu.pptx"
LOGO = ROOT / "static" / "images" / "ptt-logo.png"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT_LIGHT = RGBColor(0xBF, 0xDB, 0xFE)

STAGE_HEADERS = ["Giai đoạn", "Hành động chính", "Dữ liệu / tài liệu", "SLA"]

# fmt: off
SERVICES: list[dict] = [
    {
        "num": 1, "name": "SEO Tổng thể", "slug": "dich-vu-seo-tong-the", "group": "Tìm kiếm tự nhiên",
        "desc": "Chiến lược SEO toàn diện — kỹ thuật, nội dung, liên kết. Retainer ≥3 tháng.",
        "crm_lead": "Ngành | Ngân sách/tháng | Domain | Nhu cầu cụ thể",
        "profile": "DN có website, muốn tăng organic, giảm phụ thuộc Ads; SME hoặc enterprise.",
        "kpi": "Traffic +20% (3T), +40% (6T); ≥50% từ khóa top 10; CWV đạt chuẩn",
        "stages": [
            ["1. Lead", "AI qualify tag seo-tong-the; AM phản hồi; tạo hồ sơ CRM",
             "L0: liên hệ, nguồn lead | L1: ngành, domain, budget, pain point", "≤2h"],
            ["2. Tư vấn", "Meeting khám phá; crawl sơ bộ; so sánh 2–3 đối thủ",
             "URL; screenshot GSC/GA4; từ khóa mục tiêu; lịch sử SEO; đối thủ", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: on-page/technical/content/link; milestone 3–6T",
             "KPI cam kết; timeline; ngân sách; phạm vi IN/OUT (BRD nháp)", "≤2 ngày sau tư vấn"],
            ["4. Onboard", "Thu access; kickoff; xác nhận cluster từ khóa",
             "GSC + GA4 (Editor); từ khóa ưu tiên; brand guideline; đối thủ bám", "48–72h sau ký"],
            ["5. Triển khai", "T1 technical fix | T2 on-page+content | T3+ link building",
             "CMS access (nếu fix); content duyệt; báo cáo tiến độ hàng tháng", "6 tuần audit đầu"],
            ["6. Nghiệm thu", "Báo cáo milestone: ranking, traffic, CWV trước/sau",
             "Biên bản nghiệm thu từng mốc; CSAT ≥4.2", "Theo milestone HĐ"],
            ["7. Chăm sóc", "Báo cáo tháng; alert tụt hạng; nhắc gia hạn 30 ngày",
             "GSC/GA4 duy trì; đề xuất upsell AEO/Content", "Trước ngày 5/tháng"],
        ],
        "questions": [
            "Organic traffic/tháng hiện tại?", "Lịch sử SEO & agency cũ?",
            "Top 5 từ khóa mục tiêu?", "Quy mô site (URL/index)?",
            "CMS & ai duyệt thay đổi?", "Penalty/traffic tụt gần đây?",
        ],
        "docs_lead": "URL | Screenshot GSC Performance | GA4 overview 12 tháng",
        "docs_onboard": "GSC User | GA4 Editor | Danh sách từ khóa | Brand guideline",
        "red_flags": "Không có web | Budget <3 tháng | Kỳ vọng top 1 trong 1 tháng | Từ chối duyệt content",
        "go": "Domain live ≥3T; budget retainer; cam kết ≥3 tháng; có decision maker",
    },
    {
        "num": 2, "name": "AEO", "slug": "dich-vu-aeo", "group": "Tìm kiếm tự nhiên",
        "desc": "Tối ưu xuất hiện trong ChatGPT, Gemini, Perplexity, Google SGE.",
        "crm_lead": "Ngành | Domain | Ngân sách/tháng | Nhu cầu",
        "profile": "Đã có SEO/content; ngành cạnh tranh (BĐS, TC, YTE, giáo dục); có website nội dung.",
        "kpi": "URL được SGE/AI trích dẫn tăng; ≥80% trang ưu tiên có FAQ; 0 lỗi schema critical",
        "stages": [
            ["1. Lead", "AI tag aeo; phân loại nhu cầu AI search",
             "Ngành, website, mục tiêu sơ bộ; kiểm tra KH có web + content", "≤2h"],
            ["2. Tư vấn", "Meeting 30–45p; audit AI search presence sơ bộ",
             "URL + 5 trang ưu tiên; câu hỏi KH hay nhận; FAQ hiện có", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: chiến lược AEO, timeline, KPI xuất hiện AI",
             "Phạm vi URL; deliverable schema/FAQ; budget; BRD", "≤2 ngày"],
            ["4. Onboard", "Kickoff; checklist AEO; xác nhận mục tiêu & kênh liên lạc",
             "GSC, GA4, brand guideline, từ khóa; CMS deploy schema", "48–72h"],
            ["5. Triển khai", "T1–2 kiểm kê content | T3–4 FAQ/schema/HowTo",
             "Nội dung gốc KH; expert review (YTE/TC); duyệt ≤2 vòng", "4 tuần"],
            ["6. Nghiệm thu", "Báo cáo URL tối ưu, schema, vs KPI",
             "Rich Results Test; biên bản nghiệm thu; CSAT ≥4", "Cuối tháng 1"],
            ["7. Chăm sóc", "Báo cáo SGE/AI citation hàng tháng; alert -20%",
             "Theo dõi changelog AI providers; upsell SEO tổng thể", "Trước ngày 5"],
        ],
        "questions": [
            "Đã thử hỏi AI về thương hiệu?", "FAQ/schema hiện có?",
            "E-E-A-T: author, credentials?", "URL ưu tiên?",
            "Compliance duyệt nội dung (YTE)?", "Song song SEO tổng thể?",
        ],
        "docs_lead": "URL | Export câu hỏi CS/sales | Screenshot Rich Results",
        "docs_onboard": "GSC | GA4 | Brand guideline | CMS admin (deploy schema)",
        "red_flags": "Site mỏng nội dung | Không expert review ngành regulated | Kỳ vọng AI luôn recommend",
        "go": "≥20 URL nội dung; sẵn sàng FAQ + schema; có CMS access sau ký",
    },
    {
        "num": 3, "name": "SEO Local", "slug": "dich-vu-seo-local", "group": "Tìm kiếm tự nhiên",
        "desc": "Google Maps, GBP, tìm kiếm địa phương — Local Pack top 3.",
        "crm_lead": "Ngành | Thành phố/khu vực | Tình trạng GBP | Ngân sách/tháng",
        "profile": "Cửa hàng/chi nhánh; nhà hàng, spa, PK, BĐS, bán lẻ; cần hiển thị gần tôi.",
        "kpi": "Local Pack top 3 từ khóa ưu tiên; NAP consistency; review growth",
        "stages": [
            ["1. Lead", "AI tag seo-local; phát hiện số chi nhánh",
             "Ngành, khu vực, link GBP (nếu có), pain point Maps", "≤2h"],
            ["2. Tư vấn", "GBP audit sơ bộ; so sánh đối thủ local; 3 hành động ưu tiên",
             "Link GBP; NAP từng chi nhánh; số review; mục tiêu (gọi/chỉ đường/form)", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: setup GBP, citation, local content, review plan",
             "Số chi nhánh; KPI Local Pack %; timeline; budget/tháng", "Ngày 3–7"],
            ["4. Onboard", "Access GBP owner; xác nhận từ khóa địa phương",
             "GBP owner | Brand | Excel NAP chi nhánh | Ảnh cửa hàng | GPKD (verify)", "48–72h"],
            ["5. Triển khai", "T1 tối ưu GBP | T2 NAP+citation+local content | Hàng tháng post/review",
             "Ảnh chất lượng; mô tả GBP; website NAP khớp", "2 tuần setup"],
            ["6. Nghiệm thu", "Báo cáo GBP status, NAP score, từ khóa map",
             "Biên bản nghiệm thu setup", "Cuối tuần 2"],
            ["7. Chăm sóc", "Báo cáo view/call/chỉ đường/review mới",
             "Monitor review bombing; nhắc gia hạn", "Hàng tháng"],
        ],
        "questions": [
            "Số chi nhánh & NAP?", "GBP verify & owner?", "Review/rating?",
            "Duplicate/suspended?", "Website NAP khớp?", "Quy trình xin review?",
        ],
        "docs_lead": "Link GBP | File Excel NAP | Ảnh storefront | Bản đồ khu vực target",
        "docs_onboard": "GBP Owner | Ảnh team/SP | GPKD | Brand guideline",
        "red_flags": "GBP suspended | NAP lệch nhiều nơi | N chi nhánh, budget 1 điểm",
        "go": "Có GBP hoặc cam kết tạo; NAP cơ bản; budget theo số chi nhánh",
    },
    {
        "num": 4, "name": "SEO Audit", "slug": "dich-vu-seo-audit", "group": "Tìm kiếm tự nhiên",
        "desc": "Audit technical, on-page, content, backlink — báo cáo ưu tiên hành động.",
        "crm_lead": "Domain audit | Ngành | Ngân sách audit | Mục tiêu audit",
        "profile": "Traffic tụt, penalty nghi ngờ, redesign, mua site, cần review trước SEO dài hạn.",
        "kpi": "Báo cáo đầy đủ 2–3 tuần; bảng ưu tiên impact×effort; 0 scope creep ngoài HĐ",
        "stages": [
            ["1. Lead", "AI tag seo-audit; phát hiện urgency (traffic tụt)",
             "Domain, trigger audit, deadline báo cáo, ngân sách", "≤2h"],
            ["2. Tư vấn", "Crawl sơ bộ; 5–10 lỗi nhanh; xác định phạm vi audit",
             "URL; lịch sử migration; quy mô site; phạm vi (tech/content/link)", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: danh mục kiểm tra, tool, format, timeline",
             "Số trang giới hạn; deliverables; phí; deadline", "Ngày 3–5"],
            ["4. Onboard", "Kickoff; checklist audit — BẮT BUỘC GSC",
             "GSC + GA4 (24h) | URL ưu tiên | CMS read (nếu cần)", "24h sau ký"],
            ["5. Triển khai", "T1 technical | T2 on-page+content | T3 backlink+tổng hợp",
             "Access duy trì; stakeholder nhận báo cáo", "2–3 tuần"],
            ["6. Nghiệm thu", "Present 60p; báo cáo + bảng ưu tiên + phụ lục",
             "Biên bản nghiệm thu; upsell SEO tổng thể", "Cuối tuần 3"],
            ["7. Chăm sóc", "Follow-up implement; đề xuất triển khai theo bảng ưu tiên",
             "Tracking fix progress (nếu thuê thêm)", "30 ngày sau"],
        ],
        "questions": [
            "Trigger: tụt traffic/penalty/redesign?", "Quy mô site (URL)?",
            "Phạm vi audit mong muốn?", "Deadline báo cáo?", "Dev in-house fix?",
            "Lịch sử đổi domain/migration?",
        ],
        "docs_lead": "URL | Screenshot GSC | Báo cáo SEO cũ | Mô tả vấn đề",
        "docs_onboard": "GSC + GA4 (bắt buộc 24h) | URL ưu tiên | CMS read",
        "red_flags": "Từ chối GSC | Site >10k URL, budget nhỏ | Urgency cao, không có stakeholder",
        "go": "Cam kết GSC 24h; phạm vi & số trang thống nhất",
    },
    {
        "num": 5, "name": "Quản trị Website", "slug": "dich-vu-quan-tri-website", "group": "Tìm kiếm tự nhiên",
        "desc": "Bảo trì WordPress/custom — update, bảo mật, content, fix lỗi, backup.",
        "crm_lead": "Domain | Platform (WP/custom) | Ngành | Ngân sách/tháng",
        "profile": "Site đang vận hành; cần duy trì; vừa bàn giao thiết kế; site đang lỗi (urgency).",
        "kpi": "Uptime SLA; response time; số tác vụ hoàn thành/tháng",
        "stages": [
            ["1. Lead", "AI tag quan-tri-website; phát hiện urgency site lỗi",
             "URL, platform, mô tả lỗi, ngân sách, tần suất cập nhật", "≤2h (down ≤30p)"],
            ["2. Tư vấn", "Kiểm tra sơ bộ: tốc độ, CMS version, lỗ hổng",
             "Screenshot lỗi; hosting; plugin list; nhu cầu đặc thù", "Ngày 1–2"],
            ["3. Báo giá", "Proposal: scope/tháng, SLA uptime, số lần cập nhật",
             "Checklist bàn giao tiếp nhận; phí/tháng", "Ngày 2–5"],
            ["4. Onboard", "Hồ sơ kỹ thuật; backup đầu tiên; kickoff",
             "Hosting | wp-admin/FTP | domain | GA4 | handover dev cũ", "48–72h"],
            ["5. Triển khai", "Hàng tuần: uptime, backup, update staging trước",
             "Brief cập nhật content; ticket yêu cầu KH", "Liên tục"],
            ["6. Báo cáo", "Báo cáo tháng: uptime, tác vụ, lỗi & xử lý",
             "Incident log; đề xuất nâng cấp", "Trước ngày 5"],
            ["7. Chăm sóc", "Alert uptime <99%; gia hạn; upsell SEO/design",
             "Training KH quy tắc sửa an toàn", "Hàng tháng"],
        ],
        "questions": [
            "Lỗi hiện tại (urgency)?", "WP version & plugins?", "Hosting ai sở hữu?",
            "Tần suất cập nhật?", "Ai tự sửa admin?", "Backup & staging?",
            "SLA mong muốn?",
        ],
        "docs_lead": "URL | Mô tả lỗi + screenshot | Hosting provider",
        "docs_onboard": "wp-admin | Hosting/FTP | Domain | Handover doc",
        "red_flags": "KH tự admin hay phá | Không backup | Hosting yếu | Fix ngay chưa ký SLA",
        "go": "Site live; budget retainer; cam kết cung cấp access sau ký",
    },
    {
        "num": 6, "name": "Thiết kế Website", "slug": "thiet-ke-website", "group": "Thiết kế",
        "desc": "UI/UX design — Figma/PSD, wireframe → UI, handoff developer. 5–7 tuần.",
        "crm_lead": "Ngành | Loại web | Ngân sách | Deadline",
        "profile": "Làm mới/redesign; cần design chuyên nghiệp; dev có thể bên thứ 3.",
        "kpi": "Deliverable đúng milestone; ≤2 vòng revision miễn phí; CSAT ≥4.2",
        "stages": [
            ["1. Lead", "AI tag thiet-ke-website; redesign vs mới",
             "Loại web, số trang ước tính, budget, deadline, refs", "≤2h"],
            ["2. Tư vấn", "Meeting: brand, refs, mục tiêu CVR; benchmark 3 site",
             "Brand guideline | 3–5 URL tham khảo | sitemap sơ bộ", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: phạm vi trang, revision, timeline, format bàn giao",
             "Deliverables Figma; milestone thanh toán", "Ngày 3–7"],
            ["4. Onboard", "Design brief; kickoff xác nhận sitemap & milestone",
             "Logo SVG/EPS | Màu HEX | Font | Ảnh min 2MB | Nội dung trang", "48–72h"],
            ["5. Triển khai", "T1–2 wireframe | T3–5 UI responsive | T6–7 handoff+review",
             "Feedback tổng hợp; approval form trước code", "5–7 tuần"],
            ["6. Nghiệm thu", "Bàn giao Figma + style guide + ghi chú dev",
             "Checklist deliverable; biên bản nghiệm thu", "Cuối tuần 7"],
            ["7. Chăm sóc", "Upsell trọn gói hoặc quản trị web",
             "Hỗ trợ dev partner nếu cần", "30 ngày"],
        ],
        "questions": [
            "Redesign hay mới?", "Số trang/template?", "Refs design?",
            "Brand guideline & logo vector?", "Mục tiêu CVR?", "Số vòng revision?",
            "Dev partner có không?",
        ],
        "docs_lead": "Refs URL | Logo (nếu có) | Brand cơ bản",
        "docs_onboard": "Logo SVG/EPS | HEX colors | Font license | Ảnh ≥2MB | Sitemap/nội dung",
        "red_flags": "Không logo vector | Deadline <4T cho >10 trang | Nhiều stakeholder, không ai duyệt",
        "go": "Budget & deadline realistic; có refs; decision maker duyệt design",
    },
    {
        "num": 7, "name": "Website Trọn gói", "slug": "thiet-ke-website-tron-goi", "group": "Thiết kế",
        "desc": "Design + dev + go-live — WordPress/custom. 8–12 tuần.",
        "crm_lead": "Ngành | Loại web | Tính năng | Ngân sách",
        "profile": "Một đơn vị làm hết; landing/brochure/e-commerce; cần tích hợp payment/CRM.",
        "kpi": "Go-live đúng milestone; QA checklist; warranty bug 30 ngày",
        "stages": [
            ["1. Lead", "AI phân loại quy mô: landing/brochure/e-commerce",
             "Tính năng sơ bộ, budget, timeline, ngành", "≤2h"],
            ["2. Tư vấn", "Meeting: tính năng, tích hợp, hosting, stack",
             "BRD nháp: sitemap, user flow, feature list", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: scope design+dev, milestone, warranty",
             "Scope document IN/OUT ký trước dev; thanh toán theo mốc", "Ngày 3–7"],
            ["4. Onboard", "Brief tổng hợp; kickoff sitemap cuối, tech stack",
             "Brand | Domain/hosting | Content matrix | Integration spec", "48–72h"],
            ["5. Triển khai", "T1–2 wireframe | T3–5 UI | T6–8 dev | T9–10 QA+UAT",
             "Content deadline HĐ; staging test; tracking+SSL", "8–12 tuần"],
            ["6. Nghiệm thu", "Go-live readiness; biên bản; bàn giao access",
             "PageSpeed, mobile, form, tracking checklist", "Go-live"],
            ["7. Chăm sóc", "Warranty 30 ngày; upsell quản trị/SEO/Ads",
             "Training CMS cơ bản", "30 ngày post-live"],
        ],
        "questions": [
            "Loại web & tính năng (cart, payment, đa ngôn ngữ)?", "Tích hợp CRM/ERP?",
            "WP vs custom?", "Domain/hosting?", "Ai cung cấp content — deadline?",
            "SEO trong scope?", "Quy mô SKU (e-comm)?",
        ],
        "docs_lead": "Refs site | Feature list nháp | Brand cơ bản",
        "docs_onboard": "Scope doc ký | Brand assets | Content matrix | Domain/hosting access",
        "red_flags": "Feature creep | Content KH trễ | Không scope document | Thiếu payment spec",
        "go": "Scope ký trước dev; budget 8–12T; content plan có deadline",
    },
    {
        "num": 8, "name": "Landing Page", "slug": "thiet-ke-landing-page", "group": "Thiết kế",
        "desc": "LP chuyển đổi cao cho campaign Ads/email. 1–2 tuần.",
        "crm_lead": "Ngành | Mục đích LP | Campaign đi kèm | Ngân sách",
        "profile": "Campaign Ads sắp chạy; deadline gấp; cần CVR-focused 1 trang.",
        "kpi": "PageSpeed pass; form+tracking OK; live đúng deadline campaign",
        "stages": [
            ["1. Lead", "AI tag landing-page; urgency deadline campaign",
             "Mục đích (lead/sale/event), campaign, budget, deadline", "≤2h (≤5 ngày ≤1h)"],
            ["2. Tư vấn", "Meeting 30p: CTA, USP, section, form fields",
             "Ads creative | offer | refs LP ngành", "Ngày 1"],
            ["3. Báo giá", "Proposal gọn: scope, timeline, giá, wireframe text",
             "Số section; revision; code hay design only", "Ngày 1–2"],
            ["4. Onboard", "Brief design+copy; checklist tracking",
             "Copy USP final | Ảnh 1000×1000 | Logo SVG | Pixel/GTM ID", "Ngày 1 sau ký"],
            ["5. Triển khai", "T1–3 design | T4–7 code+form+pixel | T8–10 review",
             "Message match Ads; mobile QA", "1–2 tuần"],
            ["6. Nghiệm thu", "PageSpeed, form test, tracking events",
             "URL live; access CMS/code; biên bản", "Cuối tuần 2"],
            ["7. Chăm sóc", "Review CVR 30 ngày; upsell A/B test, Ads",
             "Hotjar/heatmap đề xuất", "30 ngày"],
        ],
        "questions": [
            "Deadline campaign?", "CTA: form/gọi/mua?", "Offer/USP/pricing?",
            "Ads creative message?", "Pixel/events?", "Domain/subdomain?",
            "Form fields & CRM nhận lead?",
        ],
        "docs_lead": "Ads creative | USP/offer | Ảnh SP | Refs LP",
        "docs_onboard": "Copy final | Logo SVG | Ảnh 1000×1000 | GTM/pixel ID",
        "red_flags": "Copy chưa chốt | Deadline <5 ngày thiếu asset | Offer yếu, kỳ vọng CVR cao",
        "go": "Deadline ≥7 ngày (ideal); copy+ảnh sẵn; campaign date xác nhận",
    },
    {
        "num": 9, "name": "Facebook Ads", "slug": "quang-cao-facebook", "group": "Quảng cáo",
        "desc": "Quản lý Meta Ads — lead gen, sales, traffic. Retainer tháng.",
        "crm_lead": "Ngành | Ngân sách/ngày | Mục tiêu campaign | Có TK Ads?",
        "profile": "Chạy Ads mới hoặc tối ưu account hiện tại; cần CPL/ROAS ổn định.",
        "kpi": "CTR/CPL/ROAS theo cam kết; báo cáo tuần+tháng đúng hạn",
        "stages": [
            ["1. Lead", "AI tag quang-cao-facebook; phân loại lead/sale/traffic",
             "Ngân sách/ngày, mục tiêu, có BM/Ads account", "≤2h"],
            ["2. Tư vấn", "Phân tích account (nếu có); benchmark CPL; strategy sơ bộ",
             "Fanpage link | screenshot Ads Manager | LP URL | USP", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: cấu trúc campaign, KPI, phí quản lý",
             "KPI CTR/CPL/ROAS; timeline; % spend nếu có", "Ngày 3–5"],
            ["4. Onboard", "Setup BM, pixel, audience; media plan T1",
             "BM partner | Pixel verify 24h | Creative kit | LP access", "48–72h"],
            ["5. Triển khai", "T1 setup campaign | T2+ daily optimize, A/B creative",
             "Creative bank 4–6 tuần; báo cáo tuần", "3–5 ngày setup"],
            ["6. Báo cáo", "Tuần: reach/CTR/CPL | Tháng: vs KPI, kế hoạch T+1",
             "Dashboard; nhận xét AM", "Tuần + trước ngày 5"],
            ["7. Chăm sóc", "Gia hạn; upsell Google Ads, retargeting",
             "Scale budget; refresh creative 2 tuần/lần", "Hàng tháng"],
        ],
        "questions": [
            "Mục tiêu lead/sale/traffic?", "KPI CPL/ROAS baseline?", "LP URL & tốc độ?",
            "Fanpage/BM owner?", "Pixel/CAPI?", "Creative ai làm?",
            "Lịch sử policy/vi phạm?",
        ],
        "docs_lead": "Fanpage | Screenshot Ads Manager | LP URL | Media kit",
        "docs_onboard": "BM partner | Pixel | Creative assets | LP/CMS",
        "red_flags": "Không LP | Pixel lỗi | Ngành hạn chế | Ngân sách/ngày quá thấp",
        "go": "LP live; spend/ngày đủ learning; BM access sau ký",
    },
    {
        "num": 10, "name": "Google Ads", "slug": "quang-cao-google", "group": "Quảng cáo",
        "desc": "Search, Shopping, Display, PMax — tối ưu CPA/ROAS.",
        "crm_lead": "Ngành | Ngân sách/tháng | Loại campaign | Có Google Ads?",
        "profile": "Search intent cao; e-commerce Shopping; ROAS/CPA target rõ.",
        "kpi": "Impression share ≥60%; CPA/ROAS cam kết; conversion data chính xác",
        "stages": [
            ["1. Lead", "AI phân loại search/shopping/display/PMax",
             "Ngân sách/tháng, mục tiêu, account hiện trạng", "≤2h"],
            ["2. Tư vấn", "Keyword research; account audit; forecast CPL/ROAS",
             "LP URL | SKU list | screenshot account | từ khóa brand", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: account structure, KPI, phí quản lý",
             "Campaign plan; budget phân bổ", "Ngày 3–5"],
            ["4. Onboard", "Conversion tracking verify; campaign structure",
             "Ads partner | GA4 | GTM | Merchant Center+feed (Shopping)", "48–72h"],
            ["5. Triển khai", "T1 setup RSA/extensions | T2+ bid, negatives, QS",
             "Search term report; landing page speed", "3–5 ngày setup"],
            ["6. Báo cáo", "Tuần + tháng: impressions, CPA, conversions vs KPI",
             "Auction insights; kế hoạch T+1", "Tuần + trước ngày 5"],
            ["7. Chăm sóc", "Gia hạn; upsell PMax/Shopping/FB Ads",
             "Budget alert 50/75/90%", "Hàng tháng"],
        ],
        "questions": [
            "Search/Shopping/PMax?", "Merchant Center/feed?", "Conversion tracking OK?",
            "ROAS/CPA target?", "Geo & ngôn ngữ?", "AOV/margin?",
            "Lịch sử account spend?",
        ],
        "docs_lead": "LP URL | SKU top sellers | Screenshot Google Ads",
        "docs_onboard": "Ads partner | GA4 | GTM | Product feed XML/Sheet",
        "red_flags": "Conversion tracking sai | LP chậm | Budget cạn giữa tháng",
        "go": "LP ≥3s load; tracking plan; budget/tháng đủ data",
    },
    {
        "num": 11, "name": "Thuê TK Quảng cáo", "slug": "thue-tai-khoan-quang-cao", "group": "Quảng cáo",
        "desc": "Thuê tài khoản Meta/Google/TikTok — setup 1–2 ngày.",
        "crm_lead": "Nền tảng | Ngân sách/tháng | Lý do thuê | Ngành",
        "profile": "TK bị khóa; cần chạy gấp; không tự mở TK; spend lớn.",
        "kpi": "TK active đúng hạn; minh bạch spend; 0 vi phạm do KH",
        "stages": [
            ["1. Lead", "AI phát hiện urgency (bị khóa + campaign đang chạy)",
             "Nền tảng, lý do thuê, spend/tháng, ngành, sản phẩm QC", "≤2h (khóa ≤1h)"],
            ["2. Tư vấn", "Đánh giá rủi ro policy; giải thích điều khoản 2 bên",
             "Lịch sử TK (screenshot) | mẫu creative | GPKD ngành regulated", "Ngày 1"],
            ["3. Báo giá/Ký", "HĐ: phí thuê, % spend, trách nhiệm, chấm dứt",
             "Draft HĐ sớm; KH cam kết tuân policy", "Ngày 1–2"],
            ["4. Onboard", "Cấu hình TK, thẻ TT, BM; hướng dẫn an toàn",
             "BM link | payment method | creative approved", "1–2 ngày"],
            ["5. Vận hành", "KH/SP chạy campaign; monitor policy daily",
             "Creative pre-approval PTT; spend report", "Hàng tháng"],
            ["6. Báo cáo", "Spend, phí thuê, tình trạng TK, sự cố",
             "Hóa đơn minh bạch", "Trước ngày 5"],
            ["7. Chăm sóc", "Nhắc gia hạn 30 ngày; upsell Ads management",
             "Backup account plan", "Hàng tháng"],
        ],
        "questions": [
            "Meta/Google/TikTok?", "Suspended/limit/mới?", "Sản phẩm quảng cáo?",
            "Spend/tháng?", "Ai tạo creative?", "Payment ai chịu?",
            "Lịch sử vi phạm?", "Ngành regulated?",
        ],
        "docs_lead": "Screenshot email platform | Mẫu creative | GPKD | Mô tả lịch sử TK",
        "docs_onboard": "BM access | Payment | Creative approved list | HĐ ký",
        "red_flags": "SP vi phạm policy | KH không cho review creative | Nhiều TK chết",
        "go": "Sản phẩm pass policy review; HĐ điều khoản ký; spend rõ ràng",
    },
    {
        "num": 12, "name": "Tiếp thị Nội dung", "slug": "tiep-thi-noi-dung", "group": "Nội dung",
        "desc": "Content marketing — blog, social; retainer bài/tháng.",
        "crm_lead": "Ngành | Kênh (blog/social) | Số bài/tháng | Ngân sách/tháng",
        "profile": "Không team viết nội bộ; cần traffic organic + brand; cần lịch publish.",
        "kpi": "100% bài đúng lịch; SEO on-page pass; traffic từ content tăng",
        "stages": [
            ["1. Lead", "AI tag tiep-thi-noi-dung; kênh blog/social/cả hai",
             "Ngành, tần suất, budget, mục tiêu traffic/brand", "≤2h"],
            ["2. Tư vấn", "Phân tích content hiện có; cluster chủ đề; calendar sơ bộ",
             "Link blog/social | bài mẫu | persona | tone | đối thủ content", "Ngày 1–3"],
            ["3. Báo giá", "Proposal: số bài/tháng, format, quy trình duyệt",
             "Content calendar T1 minh họa; SLA duyệt 48h", "Ngày 3–7"],
            ["4. Onboard", "Strategy brief; persona; brand voice guide",
             "Brand guideline | tone | từ khóa | chủ đề cấm | CMS publish", "48–72h"],
            ["5. Triển khai", "T1 plan | T2–3 viết+AI review | T4 duyệt+publish",
             "SEO brief/bài; ảnh minh họa; fact-check", "Hàng tháng"],
            ["6. Báo cáo", "Số bài publish; traffic; ranking; engagement",
             "Đề xuất cluster T+1", "Trước ngày 5"],
            ["7. Chăm sóc", "Gia hạn; upsell SEO tổng thể, Ads",
             "Calendar linh hoạt 70% plan + 30% trending", "Hàng tháng"],
        ],
        "questions": [
            "Kênh blog/social?", "Persona & tone?", "Nội dung cũ — URL?",
            "Ai duyệt — SLA?", "Chủ đề cấm/claim?", "Ai publish?",
            "Từ khóa ưu tiên?",
        ],
        "docs_lead": "Brand guideline | Bài mẫu | Link blog/social hiện tại",
        "docs_onboard": "Tone/style guide | Keyword list | CMS access | Persona doc",
        "red_flags": "Duyệt chậm >48h | Không brand voice | 100% AI không human review",
        "go": "Brand cơ bản có; người duyệt xác định; SLA 48h trong HĐ",
    },
]
# fmt: on


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _logo_exists() -> bool:
    return LOGO.is_file()


def _add_logo(slide, left=8.35, top=0.08, height=0.88):
    if _logo_exists():
        slide.shapes.add_picture(str(LOGO), Inches(left), Inches(top), height=Inches(height))


def _header(slide, title: str, subtitle: str = "", show_logo: bool = True):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    if show_logo:
        _add_logo(slide, left=8.35, top=0.08, height=0.88)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(7.7), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(7.7), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = ACCENT_LIGHT


def _bullets(slide, items: list[str], top=1.15, size=13, left=0.45, width=9.1, height=6.0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        lvl = 1 if item.startswith("  ") else 0
        p.level = lvl
        p.font.size = Pt(size - 1 if lvl else size)
        p.font.color.rgb = GRAY if lvl else DARK_TEXT
        p.space_after = Pt(4)


def _table_slide(prs, title: str, headers: list[str], rows: list[list[str]], subtitle: str = "",
                 col_widths: list[float] | None = None):
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, title, subtitle)
    nr, nc = len(rows) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(nr, nc, Inches(0.25), Inches(1.12), Inches(9.5), Inches(5.85))
    tbl = tbl_shape.table
    if col_widths and len(col_widths) == nc:
        total = sum(col_widths)
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(9.5 * w / total)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK_TEXT
            if i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT_BG
    return slide


def _service_overview_slide(prs, svc: dict):
    slide = prs.slides.add_slide(_blank(prs))
    _header(
        slide,
        f"{svc['num']}. {svc['name']}",
        f"{svc['group']} · {svc['slug']}",
    )
    items = [
        f"Mô tả: {svc['desc']}",
        f"Đối tượng: {svc['profile']}",
        f"CRM Lead (form): {svc['crm_lead']}",
        f"KPI cam kết (tóm tắt): {svc['kpi']}",
        "",
        "Câu hỏi qualify khi gặp KH:",
    ] + [f"  • {q}" for q in svc["questions"]]
    _bullets(slide, items, top=1.12, size=12)


def _service_lifecycle_slide(prs, svc: dict):
    _table_slide(
        prs,
        f"{svc['num']}. {svc['name']} — 7 giai đoạn vòng đời",
        STAGE_HEADERS,
        svc["stages"],
        subtitle="Service Delivery: Lead → Consult → Proposal → Onboard → Deliver → Handover → Retain",
        col_widths=[1.0, 2.2, 3.8, 0.8],
    )


def _service_checklist_slide(prs, svc: dict):
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, f"{svc['num']}. {svc['name']} — Checklist tài liệu & quyết định", "AM mang checklist này khi gặp KH")
    left_items = [
        "Tài liệu buổi Lead (KH mang/share):",
        f"  {svc['docs_lead']}",
        "",
        "Tài liệu Onboarding (sau ký HĐ):",
        f"  {svc['docs_onboard']}",
    ]
    right_items = [
        "Red flags — cân nhắc No-Go:",
        f"  {svc['red_flags']}",
        "",
        "Tiêu chí Go → Tư vấn:",
        f"  {svc['go']}",
    ]
    _bullets(slide, left_items, top=1.15, left=0.4, width=4.5, size=11)
    _bullets(slide, right_items, top=1.15, left=5.0, width=4.5, size=11)
    # footer note
    note = slide.shapes.add_textbox(Inches(0.4), Inches(6.85), Inches(9.2), Inches(0.45))
    np = note.text_frame.paragraphs[0]
    np.text = "Hoàn task Lead trên CRM → đồng bộ Lead care stage: first_contact → chuyển Consult (sequential)"
    np.font.size = Pt(9)
    np.font.color.rgb = GRAY
    np.font.italic = True


def _title_slide(prs: Presentation):
    slide = prs.slides.add_slide(_blank(prs))
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    if _logo_exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.1), Inches(0.85), height=Inches(1.35))
    accent = slide.shapes.add_shape(1, 0, Inches(2.35), Inches(10), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    for y, text, sz, bold, color in [
        (2.65, "Checklist Tiếp Nhận Lead", 36, True, WHITE),
        (3.35, "12 dịch vụ PTT — Dữ liệu & tài liệu theo từng bước", 18, False, ACCENT_LIGHT),
        (4.05, "Service Delivery · Giai đoạn Lead & toàn vòng đời KH", 14, False, GRAY),
        (4.55, "Phiên bản 2026-06-30 | PTT Advertising Solutions", 12, False, GRAY),
    ]:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER


def _closing_slide(prs: Presentation):
    slide = prs.slides.add_slide(_blank(prs))
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    if _logo_exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.25), Inches(2.0), height=Inches(1.2))
    for y, text, sz in [(3.45, "Cảm ơn", 40), (4.35, "PTT Advertising Solutions", 16), (4.85, "CRM Service Delivery · Creative Martech", 12)]:
        tb = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = sz >= 20
        p.font.color.rgb = WHITE if sz >= 20 else GRAY
        p.alignment = PP_ALIGN.CENTER


def build() -> Path:
    if not _logo_exists():
        raise FileNotFoundError(f"Logo not found: {LOGO}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)

    _bullets_slide = lambda prs, t, b, sub="": (
        (s := prs.slides.add_slide(_blank(prs)),
         _header(s, t, sub),
         _bullets(s, b))[0]
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Mục lục")
    toc = ["I. Khung phân tích chung (BANT+, Go/No-Go, tài sản số)"]
    for g in ("Tìm kiếm tự nhiên", "Thiết kế", "Quảng cáo", "Nội dung"):
        nums = [str(s["num"]) for s in SERVICES if s["group"] == g]
        toc.append(f"II. {g}: dịch vụ {', '.join(nums)} — 3 slide/dịch vụ")
    toc += ["III. Quy trình AM & mẫu BRD"]
    _bullets(slide, toc, top=1.15, size=13)

    _table_slide(
        prs, "Khung 4 lớp dữ liệu",
        ["Lớp", "Tên", "Thu khi nào", "Ví dụ"],
        [
            ["L0", "Identity", "Ngay có lead", "Tên, SĐT, nguồn, dịch vụ"],
            ["L1", "Qualify BANT+", "Gặp Lead", "Budget, authority, need, timeline"],
            ["L2", "Discovery", "Lead / Tư vấn", "Screenshot GSC, refs design"],
            ["L3", "BRD / file NV", "Sau tư vấn / sau ký", "Scope IN/OUT, KPI, access list"],
        ],
    )

    _table_slide(
        prs, "Go / No-Go / Nurture",
        ["Kết quả", "Điều kiện", "Hành động"],
        [
            ["Go", "Budget ≥ floor; need rõ; có authority", "Hoàn Lead → Consult 30–45p"],
            ["Nurture", "Need rõ, budget/timeline chưa sẵn sàng", "Follow-up 7–14 ngày"],
            ["No-Go", "Không fit ngành/ngân sách", "Từ chối + gợi ý dịch vụ khác"],
            ["Escalate", "HĐ lớn / enterprise", "DIR tham gia tư vấn"],
        ],
    )

    current_group = None
    for svc in SERVICES:
        if svc["group"] != current_group:
            current_group = svc["group"]
            slide = prs.slides.add_slide(_blank(prs))
            _header(slide, f"Nhóm: {current_group}", "Chi tiết từng dịch vụ — Overview · 7 bước · Checklist")
            bar = slide.shapes.add_shape(1, Inches(0.4), Inches(2.5), Inches(9.2), Inches(3.5))
            bar.fill.solid()
            bar.fill.fore_color.rgb = LIGHT_BG
            bar.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(8.6), Inches(2.5))
            names = [f"{s['num']}. {s['name']}" for s in SERVICES if s["group"] == current_group]
            p = tb.text_frame.paragraphs[0]
            p.text = " · ".join(names)
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = NAVY
            p.alignment = PP_ALIGN.CENTER

        _service_overview_slide(prs, svc)
        _service_lifecycle_slide(prs, svc)
        _service_checklist_slide(prs, svc)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Quy trình AM — Lead → Consult")
    _bullets(slide, [
        "1. Lead vào → ≤2h: phản hồi + CRM Lead + Service Lifecycle",
        "2. AI qualify → Gọi/gặp: form Lead task + BANT+",
        "3. Go/Nurture/No-Go → Go: lịch Consult + gửi checklist tài liệu",
        "4. Hoàn task Lead (sequential) → chuyển Consult",
        "5. Consult: thu L2 + BRD nháp → ≤2 ngày proposal",
        "",
        "Sau mỗi buổi gặp Lead — ghi CRM:",
        "  Decision maker | Budget | Timeline | URL/GBP/Ads | Red flags",
        "  Tài liệu KH hứa gửi + deadline | Lịch Consult | Upsell",
    ], top=1.15, size=13)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Mẫu BRD nháp (1 trang)", "Gửi KH sau buổi Lead — hoàn thiện trước Kickoff")
    _bullets(slide, [
        "1. Mục tiêu kinh doanh (SMART)",
        "2. Pain point hiện tại",
        "3. KPI: baseline → target (tool đo)",
        "4. Scope IN / OUT",
        "5. Budget & timeline",
        "6. Stakeholder (duyệt content / Ads / kỹ thuật)",
        "7. Tài sản hiện có (URL, GBP, Ads, brand)",
        "8. Rủi ro & ràng buộc (policy, penalty, deadline campaign)",
        "9. Bước tiếp theo & tài liệu KH gửi trước [ngày]",
    ], top=1.15, size=14)

    _closing_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    n = len(Presentation(str(path)).slides)
    print(f"Created: {path} ({n} slides)")
